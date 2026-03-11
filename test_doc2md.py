#!/usr/bin/env python3
"""
test_doc2md.py - Unit tests for doc2md.py

Tests cover:
- detect_format: extension parsing for supported/unsupported types
- extract_doc_metadata: metadata extraction for DOCX, PPTX, XLSX
- convert_document: markitdown integration (mocked)
- doc_to_markdown: frontmatter + heading + content assembly
- doc_to_text: plain-text output
- main CLI: argument validation, missing file, unsupported format
"""

import sys
import unittest
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

sys.path.insert(0, str(Path(__file__).parent))

import doc2md
from doc2md import (
    detect_format,
    extract_doc_metadata,
    convert_document,
    doc_to_markdown,
    doc_to_text,
    SUPPORTED_FORMATS,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_docx_bytes(title="Test Doc", author="Test Author") -> bytes:
    """Create a minimal DOCX in memory and return bytes."""
    from docx import Document
    d = Document()
    d.core_properties.title = title
    d.core_properties.author = author
    d.add_paragraph("Hello World")
    bio = BytesIO()
    d.save(bio)
    return bio.getvalue()


def _make_pptx_bytes(title="Test Slides", author="Slide Author", slides=2) -> bytes:
    """Create a minimal PPTX in memory and return bytes."""
    from pptx import Presentation
    prs = Presentation()
    prs.core_properties.title = title
    prs.core_properties.author = author
    layout = prs.slide_layouts[5]  # blank layout
    for _ in range(slides):
        prs.slides.add_slide(layout)
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _make_xlsx_bytes(title="Test Sheet", creator="Sheet Author") -> bytes:
    """Create a minimal XLSX in memory and return bytes."""
    import openpyxl
    wb = openpyxl.Workbook()
    wb.properties.title = title
    wb.properties.creator = creator
    ws = wb.active
    ws["A1"] = "Hello"
    bio = BytesIO()
    wb.save(bio)
    return bio.getvalue()


# ---------------------------------------------------------------------------
# detect_format
# ---------------------------------------------------------------------------

class TestDetectFormat(unittest.TestCase):
    def test_docx(self):
        self.assertEqual(detect_format(Path("report.docx")), "docx")

    def test_pptx(self):
        self.assertEqual(detect_format(Path("slides.pptx")), "pptx")

    def test_xlsx(self):
        self.assertEqual(detect_format(Path("data.xlsx")), "xlsx")

    def test_epub(self):
        self.assertEqual(detect_format(Path("book.epub")), "epub")

    def test_uppercase_extension(self):
        self.assertEqual(detect_format(Path("REPORT.DOCX")), "docx")

    def test_unsupported_extension(self):
        fmt = detect_format(Path("image.png"))
        self.assertNotIn(f".{fmt}", SUPPORTED_FORMATS)

    def test_supported_formats_set(self):
        for ext in [".docx", ".pptx", ".xlsx", ".epub"]:
            self.assertIn(ext, SUPPORTED_FORMATS)


# ---------------------------------------------------------------------------
# extract_doc_metadata — DOCX
# ---------------------------------------------------------------------------

class TestExtractDocxMetadata(unittest.TestCase):
    def setUp(self):
        import tempfile, os
        self.docx_bytes = _make_docx_bytes(title="My Report", author="Jane Doe")
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
            f.write(self.docx_bytes)
            self.docx_path = Path(f.name)

    def tearDown(self):
        self.docx_path.unlink(missing_ok=True)

    def test_title_extracted(self):
        meta = extract_doc_metadata(self.docx_path, "docx")
        self.assertEqual(meta.get("title"), "My Report")

    def test_author_extracted(self):
        meta = extract_doc_metadata(self.docx_path, "docx")
        self.assertEqual(meta.get("author"), "Jane Doe")

    def test_format_field_present(self):
        meta = extract_doc_metadata(self.docx_path, "docx")
        self.assertEqual(meta.get("format"), "docx")

    def test_source_field_is_absolute(self):
        meta = extract_doc_metadata(self.docx_path, "docx")
        self.assertTrue(Path(meta["source"]).is_absolute())

    def test_fetched_at_is_iso8601(self):
        meta = extract_doc_metadata(self.docx_path, "docx")
        fetched = meta.get("fetched_at", "")
        # Should parse without error
        datetime.strptime(fetched, "%Y-%m-%dT%H:%M:%SZ")


# ---------------------------------------------------------------------------
# extract_doc_metadata — PPTX
# ---------------------------------------------------------------------------

class TestExtractPptxMetadata(unittest.TestCase):
    def setUp(self):
        import tempfile
        self.pptx_bytes = _make_pptx_bytes(title="Q4 Results", author="Bob Smith", slides=3)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(self.pptx_bytes)
            self.pptx_path = Path(f.name)

    def tearDown(self):
        self.pptx_path.unlink(missing_ok=True)

    def test_title_extracted(self):
        meta = extract_doc_metadata(self.pptx_path, "pptx")
        self.assertEqual(meta.get("title"), "Q4 Results")

    def test_author_extracted(self):
        meta = extract_doc_metadata(self.pptx_path, "pptx")
        self.assertEqual(meta.get("author"), "Bob Smith")

    def test_slides_count(self):
        meta = extract_doc_metadata(self.pptx_path, "pptx")
        self.assertEqual(meta.get("slides"), 3)

    def test_format_field(self):
        meta = extract_doc_metadata(self.pptx_path, "pptx")
        self.assertEqual(meta.get("format"), "pptx")


# ---------------------------------------------------------------------------
# extract_doc_metadata — XLSX
# ---------------------------------------------------------------------------

class TestExtractXlsxMetadata(unittest.TestCase):
    def setUp(self):
        import tempfile
        self.xlsx_bytes = _make_xlsx_bytes(title="Sales Data", creator="Alice")
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
            f.write(self.xlsx_bytes)
            self.xlsx_path = Path(f.name)

    def tearDown(self):
        self.xlsx_path.unlink(missing_ok=True)

    def test_title_extracted(self):
        meta = extract_doc_metadata(self.xlsx_path, "xlsx")
        self.assertEqual(meta.get("title"), "Sales Data")

    def test_author_extracted(self):
        meta = extract_doc_metadata(self.xlsx_path, "xlsx")
        self.assertEqual(meta.get("author"), "Alice")

    def test_sheets_count(self):
        meta = extract_doc_metadata(self.xlsx_path, "xlsx")
        self.assertGreaterEqual(meta.get("sheets", 0), 1)

    def test_format_field(self):
        meta = extract_doc_metadata(self.xlsx_path, "xlsx")
        self.assertEqual(meta.get("format"), "xlsx")


# ---------------------------------------------------------------------------
# extract_doc_metadata — fallback (epub/unknown)
# ---------------------------------------------------------------------------

class TestExtractFallbackMetadata(unittest.TestCase):
    def test_epub_has_format_and_source(self):
        # We just need a path that exists — epub metadata extractor is not implemented
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as f:
            f.write(b"fake epub")
            epub_path = Path(f.name)
        try:
            meta = extract_doc_metadata(epub_path, "epub")
            self.assertEqual(meta.get("format"), "epub")
            self.assertIn("source", meta)
            self.assertIn("fetched_at", meta)
        finally:
            epub_path.unlink(missing_ok=True)


# ---------------------------------------------------------------------------
# convert_document (mocked)
# ---------------------------------------------------------------------------

class TestConvertDocument(unittest.TestCase):
    def test_converts_docx_returns_text(self):
        mock_result = MagicMock()
        mock_result.text_content = "# Hello\n\nThis is content."

        mock_mid_cls = MagicMock()
        mock_mid_instance = mock_mid_cls.return_value
        mock_mid_instance.convert.return_value = mock_result

        mock_markitdown_module = MagicMock()
        mock_markitdown_module.MarkItDown = mock_mid_cls

        with patch.dict("sys.modules", {"markitdown": mock_markitdown_module}):
            result = convert_document(Path("fake.docx"))

        self.assertEqual(result, "# Hello\n\nThis is content.")
        mock_mid_instance.convert.assert_called_once_with("fake.docx")

    def test_raises_import_error_if_not_installed(self):
        import builtins
        real_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "markitdown":
                raise ImportError("No module named 'markitdown'")
            return real_import(name, *args, **kwargs)

        with patch("builtins.__import__", side_effect=mock_import):
            with self.assertRaises(ImportError):
                convert_document(Path("fake.docx"))


# ---------------------------------------------------------------------------
# doc_to_markdown / doc_to_text
# ---------------------------------------------------------------------------

class TestDocToMarkdown(unittest.TestCase):
    def test_includes_frontmatter(self):
        meta = {"title": "My Doc", "format": "docx", "fetched_at": "2026-03-10T00:00:00Z"}
        result = doc_to_markdown("Some content", metadata=meta, title="My Doc")
        self.assertIn("---", result)
        self.assertIn("title:", result)
        self.assertIn("Some content", result)

    def test_includes_h1_title(self):
        result = doc_to_markdown("Body text", metadata=None, title="The Title")
        self.assertIn("# The Title", result)
        self.assertIn("Body text", result)

    def test_no_metadata_no_frontmatter(self):
        result = doc_to_markdown("Body only", metadata=None, title=None)
        self.assertNotIn("---", result)
        self.assertIn("Body only", result)

    def test_format_key_order(self):
        meta = {
            "title": "Test",
            "author": "Author",
            "format": "docx",
            "fetched_at": "2026-03-10T00:00:00Z",
        }
        result = doc_to_markdown("content", metadata=meta)
        title_pos = result.find("title:")
        author_pos = result.find("author:")
        self.assertLess(title_pos, author_pos)


class TestDocToText(unittest.TestCase):
    def test_returns_content_unchanged(self):
        content = "Plain text with **markdown** syntax."
        result = doc_to_text(content)
        self.assertEqual(result, content)

    def test_no_frontmatter_in_text(self):
        content = "Just text"
        result = doc_to_text(content)
        self.assertNotIn("---", result)


# ---------------------------------------------------------------------------
# CLI — argument validation
# ---------------------------------------------------------------------------

class TestCLI(unittest.TestCase):
    def test_help_exits_cleanly(self):
        from typer.testing import CliRunner
        from doc2md import app
        runner = CliRunner()
        result = runner.invoke(app, ["--help"])
        self.assertEqual(result.exit_code, 0)
        self.assertIn("DOCX", result.output)

    def test_missing_file_exits_with_error(self):
        from typer.testing import CliRunner
        from doc2md import app
        runner = CliRunner()
        result = runner.invoke(app, ["/nonexistent/path/file.docx"])
        self.assertNotEqual(result.exit_code, 0)

    def test_unsupported_format_exits_with_error(self):
        import tempfile
        from typer.testing import CliRunner
        from doc2md import app
        runner = CliRunner()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            f.write(b"fake image")
            png_path = Path(f.name)
        try:
            result = runner.invoke(app, [str(png_path)])
            self.assertNotEqual(result.exit_code, 0)
        finally:
            png_path.unlink(missing_ok=True)

    def test_converts_docx_end_to_end(self):
        """Integration smoke test: convert a real DOCX to markdown."""
        import tempfile
        from typer.testing import CliRunner
        from doc2md import app

        docx_bytes = _make_docx_bytes(title="E2E Test", author="Tester")
        with tempfile.TemporaryDirectory() as tmpdir:
            docx_path = Path(tmpdir) / "test.docx"
            docx_path.write_bytes(docx_bytes)
            out_dir = Path(tmpdir) / "out"

            runner = CliRunner()
            result = runner.invoke(app, [str(docx_path), "-o", str(out_dir)])

            self.assertEqual(result.exit_code, 0, msg=result.output)
            out_file = out_dir / "test.md"
            self.assertTrue(out_file.exists())
            content = out_file.read_text()
            self.assertIn("---", content)
            self.assertIn("Hello World", content)


if __name__ == "__main__":
    unittest.main()
