#!/usr/bin/env python3
"""
test/create_fixtures.py — Generate all test fixtures programmatically.

Run this script once to populate test/fixtures/ with real files
used by test/test_fixtures.py and test_integration.py.

No downloads required — everything is created from library calls.

Usage:
    python test/create_fixtures.py
"""

from pathlib import Path

FIXTURES = Path(__file__).parent / "fixtures"
FIXTURES.mkdir(parents=True, exist_ok=True)


def create_rst():
    """Write sample.rst — RST with docinfo fields, code block, and table."""
    content = """\
Test Document
=============

:Author: Test Author
:Date: 2026-01-01
:Version: 1.0

Introduction
------------

This is a test RST document with **bold** and *italic* text.

.. code-block:: python

   print("hello world")

A table:

+-------+-------+
| Col 1 | Col 2 |
+-------+-------+
| a     | b     |
+-------+-------+
"""
    path = FIXTURES / "sample.rst"
    path.write_text(content, encoding="utf-8")
    print(f"Created: {path}")


def create_html():
    """Write sample.html — HTML with meta tags and body content."""
    content = """\
<!DOCTYPE html>
<html>
<head>
  <title>Test Page</title>
  <meta name="description" content="A test page for html2md">
  <meta name="author" content="Test Author">
</head>
<body>
  <h1>Main Heading</h1>
  <p>This is a paragraph with <strong>bold</strong> and <em>italic</em> text.</p>
  <h2>Section Two</h2>
  <p>Another paragraph with a <a href="https://example.com">link</a>.</p>
  <ul>
    <li>Item one</li>
    <li>Item two</li>
  </ul>
</body>
</html>
"""
    path = FIXTURES / "sample.html"
    path.write_text(content, encoding="utf-8")
    print(f"Created: {path}")


def create_jpg():
    """Create sample.jpg — minimal valid JPEG with text drawn on white background."""
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        print("WARNING: Pillow not available; skipping sample.jpg")
        return

    img = Image.new("RGB", (200, 100), color="white")
    draw = ImageDraw.Draw(img)
    draw.text((10, 40), "Hello World", fill="black")
    path = FIXTURES / "sample.jpg"
    img.save(str(path), "JPEG")
    print(f"Created: {path}")


def create_docx():
    """Create sample.docx — minimal DOCX with title, author, and two paragraphs."""
    try:
        from docx import Document
    except ImportError:
        print("WARNING: python-docx not available; skipping sample.docx")
        return

    doc = Document()
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    doc.add_heading("Test Document", 0)
    doc.add_paragraph("This is a test paragraph.")
    doc.add_paragraph("Second paragraph with more content.")
    path = FIXTURES / "sample.docx"
    doc.save(str(path))
    print(f"Created: {path}")


def create_pptx():
    """Create sample.pptx — 2-slide PPTX with title and author set."""
    try:
        from pptx import Presentation
    except ImportError:
        print("WARNING: python-pptx not available; skipping sample.pptx")
        return

    prs = Presentation()
    prs.core_properties.author = "Test Author"

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Test Presentation"

    # Slide 2: Title + content
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Slide Two"

    path = FIXTURES / "sample.pptx"
    prs.save(str(path))
    print(f"Created: {path}")


def create_xlsx():
    """Create sample.xlsx — workbook with 1 sheet, 3 data rows."""
    try:
        import openpyxl
    except ImportError:
        print("WARNING: openpyxl not available; skipping sample.xlsx")
        return

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["foo", 1])
    ws.append(["bar", 2])
    wb.properties.title = "Test Workbook"
    wb.properties.creator = "Test Author"

    path = FIXTURES / "sample.xlsx"
    wb.save(str(path))
    print(f"Created: {path}")


def create_pdf():
    """Create sample.pdf — born-digital PDF (1 page) using PyMuPDF."""
    try:
        import fitz  # PyMuPDF — already a dep via pymupdf4llm
    except ImportError:
        print("WARNING: PyMuPDF (fitz) not available; skipping sample.pdf")
        return

    doc = fitz.open()
    page = doc.new_page()
    # Insert enough text to exceed THIN_PAGE_THRESHOLD (50 chars)
    text = (
        "Test PDF Document\n\n"
        "This is a test paragraph for pdf2md testing.\n"
        "It contains enough text to pass the THIN_PAGE_THRESHOLD.\n"
        "Additional content to ensure the page is not flagged as thin.\n"
        "Line four: more content here to pad the text layer.\n"
        "Line five: yet more content so extraction tests pass cleanly."
    )
    page.insert_text((72, 72), text)
    path = FIXTURES / "sample.pdf"
    doc.save(str(path))
    doc.close()
    print(f"Created: {path}")


def create_two_speakers_wav():
    """Download test_audio/two_speakers.wav from HuggingFace VoxConverse dataset.

    Source: https://huggingface.co/datasets/diarizers-community/voxconverse
    Split: dev, shard 0, row 4 — 44.5s clip with 2 speakers, 9 segments.
    """
    import subprocess

    test_audio = Path(__file__).parent.parent / "test_audio"
    test_audio.mkdir(parents=True, exist_ok=True)
    output = test_audio / "two_speakers.wav"

    if output.exists():
        print(f"Already exists: {output}")
        return

    try:
        import pandas as pd
        from huggingface_hub import hf_hub_download
    except ImportError:
        print("WARNING: pandas/huggingface_hub not available; skipping two_speakers.wav")
        return

    parquet = hf_hub_download(
        "diarizers-community/voxconverse",
        "data/dev-00000-of-00005.parquet",
        repo_type="dataset",
    )
    df = pd.read_parquet(parquet)
    row = df.iloc[4]  # 2-speaker, ~45s, 9 segments

    raw = test_audio / "_tmp_raw.wav"
    raw.write_bytes(row["audio"]["bytes"])
    subprocess.run(
        ["ffmpeg", "-y", "-i", str(raw), "-ar", "16000", "-ac", "1", str(output)],
        capture_output=True, check=True,
    )
    raw.unlink()
    print(f"Created: {output}")


def create_yt_interview_wav():
    """Download test_audio/yt_interview.wav from HuggingFace VoxConverse dataset.

    Source: https://huggingface.co/datasets/diarizers-community/voxconverse
    Split: test, shard 0, row 3 — trimmed to 90s, 2 speakers.
    """
    import subprocess

    test_audio = Path(__file__).parent.parent / "test_audio"
    test_audio.mkdir(parents=True, exist_ok=True)
    output = test_audio / "yt_interview.wav"

    if output.exists():
        print(f"Already exists: {output}")
        return

    try:
        import pandas as pd
        from huggingface_hub import hf_hub_download
    except ImportError:
        print("WARNING: pandas/huggingface_hub not available; skipping yt_interview.wav")
        return

    parquet = hf_hub_download(
        "diarizers-community/voxconverse",
        "data/test-00000-of-00011.parquet",
        repo_type="dataset",
    )
    df = pd.read_parquet(parquet)
    row = df.iloc[3]  # 2-speaker, ~572s full

    raw = test_audio / "_tmp_raw.wav"
    raw.write_bytes(row["audio"]["bytes"])
    subprocess.run(
        ["ffmpeg", "-y", "-i", str(raw), "-ar", "16000", "-ac", "1", "-t", "90", str(output)],
        capture_output=True, check=True,
    )
    raw.unlink()
    print(f"Created: {output}")


def main():
    print(f"Creating fixtures in: {FIXTURES}")
    create_rst()
    create_html()
    create_jpg()
    create_docx()
    create_pptx()
    create_xlsx()
    create_pdf()
    create_two_speakers_wav()
    create_yt_interview_wav()
    print("\nDone. Fixtures:")
    for f in sorted(FIXTURES.iterdir()):
        size = f.stat().st_size
        print(f"  {f.name}  ({size} bytes)")


if __name__ == "__main__":
    main()
